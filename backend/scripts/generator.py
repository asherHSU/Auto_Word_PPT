import sys
import json
import os
import re
import io

# 設定標準輸出為 UTF-8，避免中文亂碼
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

# 引入 docx 和 pptx 相關套件
try:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from pptx.dml.color import RGBColor as PptRGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print(json.dumps({"error": f"Missing dependency: {e}"}))
    sys.exit(1)

# --- 參數解析 ---
try:
    MODE = sys.argv[1]
    songs_input = json.loads(sys.argv[2])
    RESOURCES_DIR = sys.argv[3]
    OUTPUT_DIR = sys.argv[4] if len(sys.argv) > 4 else None
except IndexError:
    print(json.dumps({"error": "Missing arguments"}))
    sys.exit(1)

PPT_LIBRARY_PATH = os.path.join(RESOURCES_DIR, "ppt_library")
TEMPLATE_PATH = os.path.join(RESOURCES_DIR, "template.docx")

# --- 輔助函式 ---

def clean_text(text):
    """清除 XML 不支援的控制字元"""
    return re.sub(r'[\x00-\x08\x0e-\x1f]', '', text)

def normalize_string(s):
    """正規化字串：去除非英數中文並轉小寫，用於比對檔名"""
    return re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', '', s).lower()

def find_ppt_path(root_path, song_id, song_name):
    """遞迴搜尋 PPT 檔案"""
    if not os.path.exists(root_path): return None
    target_name = normalize_string(song_name)
    id_str = str(song_id)
    id_regex = re.compile(rf"^0*{id_str}([^0-9]|$)")

    for dirpath, _, filenames in os.walk(root_path):
        for filename in filenames:
            if not filename.lower().endswith(('.pptx', '.ppt')): continue
            name_stem = os.path.splitext(filename)[0]
            if id_regex.match(name_stem): return os.path.join(dirpath, filename)
            if target_name in normalize_string(name_stem): return os.path.join(dirpath, filename)
    return None

def extract_lyrics_from_ppt(ppt_path):
    """從 PPTX 提取歌詞"""
    if not ppt_path or not ppt_path.lower().endswith(".pptx"): return []
    try:
        prs = Presentation(ppt_path)
        slides_lyrics = []
        for slide in prs.slides:
            slide_lines = []
            for shape in slide.shapes:
                if not shape.has_text_frame or not shape.text.strip(): continue
                if shape.top < prs.slide_height * 0.9:
                    text = shape.text_frame.text.replace('\v', '\n')
                    for line in text.splitlines():
                        if line.strip():
                            cleaned_line = re.sub(r'[ \t]+', ' ', clean_text(line))
                            slide_lines.append(cleaned_line)
            if slide_lines:
                slides_lyrics.append(slide_lines)
        return slides_lyrics
    except:
        return []

def apply_font_settings(paragraph, font_name, font_size, font_color, is_bold=False):
    """設定 PPT 文字樣式"""
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = PptPt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = is_bold

def calculate_optimal_font_size(lines):
    """
    🔍 智慧字體大小計算 (Smart Font Sizing)
    
    同時考量：
    1. 最長單行的字元寬度 (避免水平超出)
    2. 總行數的堆疊高度 (避免垂直超出)
    
    Args:
        lines (list): 該頁面的所有行文字
    Returns:
        int: 計算出的最佳字體大小 (points)
    """
    if not lines: return 48 # 預設大字體
    
    # 1. 計算最長那行的「視覺寬度」
    max_visual_width = 0
    for line in lines:
        current_width = 0
        for char in line:
            # 簡單權重：全形字(中文)算 1，半形字(英文/數字)算 0.55
            if ord(char) > 127: 
                current_width += 1
            else:
                current_width += 0.55
        if current_width > max_visual_width:
            max_visual_width = current_width

    # 2. 定義畫布限制 (Points)
    # PPT 寬度 10吋，左右邊界各 0.5吋 -> 可用寬度 9吋 = 648 pt
    # PPT 高度 5.625吋，扣掉標題與邊界 -> 可用高度約 3.5吋 = 252 pt
    
    SAFE_WIDTH_PTS = 610  # 保險起見稍微縮小
    SAFE_HEIGHT_PTS = 230 # 用於歌詞的垂直空間
    
    # 3. 根據「寬度」計算上限
    # 假設字體大小為 S，全形字寬度約為 S
    # S * max_visual_width <= SAFE_WIDTH_PTS
    if max_visual_width < 1: max_visual_width = 1
    size_limit_by_width = int(SAFE_WIDTH_PTS / max_visual_width)
    
    # 4. 根據「高度」計算上限
    # 假設行高倍率為 1.15
    # total_lines * S * 1.15 <= SAFE_HEIGHT_PTS
    line_count = len(lines)
    if line_count < 1: line_count = 1
    size_limit_by_height = int(SAFE_HEIGHT_PTS / (line_count * 1.15))
    
    # 5. 取兩者最小值，並設定合理的上下限
    final_size = min(size_limit_by_width, size_limit_by_height, 54) # 最大不超過 54
    
    if final_size < 24: final_size = 24 # 最小不低於 24 (再小就看不到了，讓它自動換行)
    
    return final_size

def group_lyrics_dynamic(lyrics_list, max_lines=2):
    """動態歌詞分組演算法"""
    groups = []
    current_buffer = []
    markers = (
        '1', '2', '3', '4', '5', '6', '7', '8', '9', 
        'c', 'b', 'v', 'p', 't', 'e', 
        '§', '※', '©', '®', '＊', '*', 
        'bridge', 'chorus', 'verse', 'pre-chorus', 'tag', 'ending',
        '(1)', '(2)', '(3)', '(4)', '(c)', '(b)', '(v)', 
        '（1）', '（2）', '（3）', '（c）', '（b）', 
        '[1]', '[2]', '[c]', '[b]',
        'verse', 'chorus', 'pc', 'p-c'
    )

    for line in lyrics_list:
        clean_line = line.strip()
        if not clean_line: continue 

        is_marker_line = clean_line.lower().startswith(markers)
        
        if len(current_buffer) >= max_lines or (current_buffer and is_marker_line):
            groups.append(current_buffer)
            current_buffer = [] 

        current_buffer.append(line)

    if current_buffer:
        groups.append(current_buffer)
        
    return groups

# --- 主邏輯 ---

if MODE == "preview":
    results = []
    for song in songs_input:
        song_id = song.get('id', 0)
        song_name = song.get('name', '') or song.get('title', '')
        
        ppt_path = find_ppt_path(PPT_LIBRARY_PATH, song_id, song_name)
        lyrics_by_slide = extract_lyrics_from_ppt(ppt_path)
        flat_lyrics = [line for slide in lyrics_by_slide for line in slide]
        
        results.append({
            "title": song_name,
            "lyrics": flat_lyrics,
            "found": bool(ppt_path),
            "isOld": ppt_path.endswith(".ppt") if ppt_path else False
        })
    print(json.dumps(results, ensure_ascii=False))

elif MODE == "generate":
    if not OUTPUT_DIR:
        print(json.dumps({"error": "Output directory required"}))
        sys.exit(1)

    OUTPUT_DOCX = os.path.join(OUTPUT_DIR, "敬拜大字報.docx")
    OUTPUT_PPTX = os.path.join(OUTPUT_DIR, "敬拜PPT.pptx")

    # Word 初始化
    try:
        doc = Document(TEMPLATE_PATH)
        if doc.paragraphs and not doc.paragraphs[0].text.strip():
            p_element = doc.paragraphs[0]._element
            p_element.getparent().remove(p_element)
    except:
        doc = Document()

    # PPT 初始化
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    BLACK_FILL = PptRGBColor(0, 0, 0)
    YELLOW_TEXT = PptRGBColor(255, 255, 0)
    FONT_NAME = "微軟正黑體"

    for i, song in enumerate(songs_input):
        title = song.get('title') or song.get('name')
        is_from_preview = 'lyrics' in song and song['lyrics'] is not None
        
        lyrics_data = []
        if is_from_preview:
            lyrics_data = [re.sub(r'[ \t]+', ' ', line) for line in song['lyrics']]
        else:
            sid = song.get('id', 0)
            path_found = find_ppt_path(PPT_LIBRARY_PATH, sid, title)
            lyrics_data = extract_lyrics_from_ppt(path_found)

        # 生成 Word
        flat_lyrics = lyrics_data if is_from_preview else [line for slide in lyrics_data for line in slide]

        if i > 0: doc.add_paragraph("")
        doc.add_paragraph(f"【{title}】", style='SongTitle' if 'SongTitle' in doc.styles else None)
        
        if not flat_lyrics:
            p = doc.add_paragraph()
            run = p.add_run("【無歌詞內容】")
            run.font.color.rgb = RGBColor(0, 0, 0)
        else:
            for line in flat_lyrics:
                p = doc.add_paragraph(line, style='Lyrics' if 'Lyrics' in doc.styles else None)
                p.paragraph_format.space_before = Pt(0)
                p.paragraph_format.space_after = Pt(0)
                if line.strip().lower().startswith(('c.', 'b.', 'bridge', 'chorus')):
                    for run in p.runs: run.font.bold = True
        
        # 生成 PPT
        if not lyrics_data: continue
        layout = prs.slide_layouts[6] 

        # 決定如何分頁 (分組)
        final_slides_content = []
        if is_from_preview:
            final_slides_content = group_lyrics_dynamic(lyrics_data, max_lines=2)
        else:
            final_slides_content = lyrics_data # 原始檔案結構

        for slide_lines in final_slides_content:
            slide = prs.slides.add_slide(layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = BLACK_FILL

            lyric_text = "\n".join(slide_lines)
            
            # ✨✨✨ 使用新的智慧字體計算 ✨✨✨
            font_size = calculate_optimal_font_size(slide_lines)

            # 歌詞文字方塊
            tb_lyrics = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(3.5))
            tf_lyrics = tb_lyrics.text_frame
            tf_lyrics.word_wrap = True # 允許自動換行 (作為最後防線)
            
            p_lyrics = tf_lyrics.paragraphs[0]
            p_lyrics.text = lyric_text
            p_lyrics.alignment = PP_ALIGN.CENTER
            apply_font_settings(p_lyrics, FONT_NAME, font_size, YELLOW_TEXT, True)

            # Footer
            tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
            tf_title = tb_title.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = f"《{title}》"
            p_title.alignment = PP_ALIGN.CENTER
            apply_font_settings(p_title, FONT_NAME, 20, YELLOW_TEXT, False)

    doc.save(OUTPUT_DOCX)
    prs.save(OUTPUT_PPTX)
    print(json.dumps({"status": "success", "files": [OUTPUT_DOCX, OUTPUT_PPTX]}))

else:
    print(json.dumps({"error": "Unknown mode"}))