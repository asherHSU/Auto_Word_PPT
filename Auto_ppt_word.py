import os
import re
from docx import Document
# ▼▼▼ 匯入調整格式需要的新工具 ▼▼▼
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
# ▲▲▲ 匯入結束 ▲▲▲
from pptx import Presentation

# --- 使用者設定區 ---
song_order = [
    "將天敞開",
    "在這裡",
    "我們的神",
    "耶穌永遠掌權",
    "禱告的力量"
]
ppt_library_path = "D:/Desktop/church/Auto_ppt_word/2025 別是巴聖教會雲端詩歌PPT" # 我根據你的截圖更新了路徑
output_filename = "D:/Desktop/church/Auto_ppt_word/敬拜大字報.docx"

# --- 程式準備區 ---
# ▼▼▼ 修改：不再建立全新文件，而是從我們的範本載入 ▼▼▼
try:
    document = Document("D:\Desktop\church\Auto_ppt_word/template.docx")
except Exception as e:
    print("錯誤：找不到 'template.docx' 範本檔案，請確認它與Python腳本放在同一個資料夾。")
    print(f"詳細錯誤：{e}")
    exit()
# ▲▲▲ 修改結束 ▲▲▲

# --- 核心功能函式 ---
# --- 核心功能函式 ---
def find_ppt_path(root_path, song_title):
    """
    升級版搜尋：從檔名中提取出乾淨的標題，進行精準比對。
    """
    for dirpath, _, filenames in os.walk(root_path):
        for filename in filenames:
            # 先確認是PPT檔案
            if not (filename.lower().endswith(".pptx") or filename.lower().endswith(".ppt")):
                continue

            # 1. 去除副檔名，取得乾淨的檔名主幹 (例如 "1001-神正在這裡")
            name_stem = os.path.splitext(filename)[0]

            # 2. 從檔名主幹中，分離出真正的歌曲標題
            title_from_file = ""
            # 我們假設編號和歌名是用 '-' 或 ' ' (空格) 分隔
            if '-' in name_stem:
                # 從第一個 '-' 後面取所有文字作為標題
                title_from_file = name_stem.split('-', 1)[-1].strip()
            elif ' ' in name_stem:
                # 從第一個 ' ' 後面取所有文字作為標題
                title_from_file = name_stem.split(' ', 1)[-1].strip()
            else:
                # 如果檔名中沒有分隔符，就認為整個檔名都是標題
                title_from_file = name_stem.strip()
            
            # 3. 進行精準的「完全等於」比對
            if title_from_file == song_title:
                # 只有在完全相符時，才回傳路徑
                return os.path.join(dirpath, filename)

    # 如果整個迴圈跑完都沒找到完全相符的，就回傳 None
    return None

# ▼▼▼ 修正一：修正 clean_text 的過濾規則，不再移除換行符號 ▼▼▼
def clean_text(text):
    """
    清除文字中不相容XML的控制字元，但保留換行、歸位和定位符。
    """
    # 這個規則移除了大部分控制字元，但會保留 \t, \n, \r, \v
    return re.sub(r'[\x00-\x08\x0e-\x1f]', '', text)
# ▲▲▲ 修正結束 ▲▲▲

# --- 程式主體 ---
print("程式已啟動，準備開始處理...")
for song in song_order:
    print(f"正在處理歌曲: {song}...")
    ppt_path = find_ppt_path(ppt_library_path, song)
    
    if ppt_path:
        print(f"  > 找到了: {os.path.basename(ppt_path)}")
        
        if ppt_path.endswith(".ppt"):
            # ... (跳過 .ppt 的邏輯維持不變)
            document.add_heading(f"【{song}】", level=1)
            p = document.add_paragraph()
            p.add_run("【此歌曲為舊版.ppt格式，無法自動匯入，請手動處理】").bold = True
            document.add_paragraph("")
            continue

        document.add_paragraph(f"【{song}】", style='SongTitle')
        
        prs = Presentation(ppt_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame or not shape.text.strip():
                    continue

                if shape.top < prs.slide_height / 2:
                    # ▼▼▼ 修正二：簡化並統一處理所有換行 ▼▼▼
                    
                    # 1. 取得文字框的完整文字，包含所有 \n 和 \v 換行符
                    raw_text = shape.text_frame.text
                    
                    # 2. 將軟換行(\v)統一替換成標準換行(\n)
                    text_with_standard_breaks = raw_text.replace('\v', '\n')
                    
                    # 3. 使用 splitlines() 這個強大的方法來切割所有行
                    lines = text_with_standard_breaks.splitlines()

                    # 4. 遍歷所有切割出來的行
                    for line in lines:
                        if not line.strip():
                            continue
                        
                        # 清洗每一行文字
                        lyric_text = clean_text(line)
                        
                        # 為每一行歌詞建立一個新的Word段落
                        p = document.add_paragraph(lyric_text, style='Lyrics')
                        
                        # 樣式微調邏輯 (例如標示副歌)
                        if lyric_text.strip().startswith(('c.', 'b.')):
                            if p.runs:
                                run = p.runs[0]
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(255, 0, 0)
                    # ▲▲▲ 修改結束 ▲▲▲
        
    else:
        print(f"  > 警告: 在 '{ppt_library_path}' 中找不到歌曲 '{song}' 的PPT檔案。")

# --- 程式結束區 ---
document.save(output_filename)
print("--------------------")
print(f"處理完成！檔案 '{output_filename}' 已成功建立！")

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN  # Import PP_ALIGN for text alignment
from pptx.util import Inches

# Function to apply font settings
def apply_font_settings(paragraph, font_name, font_size, font_color, is_bold=False):
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = font_color
    paragraph.font.bold = is_bold
    
#要修改成大字報與PPT輸出的位置
doc = Document(output_filename)
output_path = "D:/Desktop/church/Auto_ppt_word/敬拜PPT.pptx"

# Create a PowerPoint presentation object
presentation = Presentation()

# Set slide size for 16:9 aspect ratio
presentation.slide_width = Inches(10)  # Width in inches
presentation.slide_height = Inches(5.625)  # Height in inches

# Define slide background and text colors
black_fill = RGBColor(0, 0, 0)
yellow_text = RGBColor(255, 255, 0)
font_name = "微軟正黑體"  # Default font name
font_size_main = 32  # Font size for main text
font_size_bottom = 20  # Font size for bottom text
p_text = ""  # To store the last encountered P line

# List to collect text chunks
text_chunks = []
p_text = []
p_text_Num = -1

# Iterate through each paragraph in the Word document
for para in doc.paragraphs:
    # Skip empty paragraphs
    if not para.text.strip():
        continue
    
    # Check if the paragraph contains '【'
    if '【' in para.text.strip():
        p_text.append(para.text.strip())  # Update the ( text
        text_chunks.append(para.text.strip())
        # continue  # Skip adding this line to the slides
    
    # Split paragraph into chunks
    words = para.text.strip().split(' ')
    max_words_per_line = 12  # Roughly one line of text
    
    # Split the text into chunks that represent one line
    for i in range(0, len(words), max_words_per_line):
        line_text = ' '.join(words[i:i + max_words_per_line])
        text_chunks.append(line_text)

p_text.append(" ")

# Add two lines of text per slide
for i in range(0, len(text_chunks), 2):
    
    if text_chunks[i] == p_text[p_text_Num + 1]:
        p_text_Num += 1
        continue

    # Add a new slide with a blank layout
    slide_layout = presentation.slide_layouts[5]  # Blank slide layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Set slide background to black
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = black_fill
    
    # Add a text box for the first line of text
    left = Inches(0.5)
    top = Inches(0.1)  # Top of the slide
    width = Inches(9)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text_chunks[i]  # First line of text
    apply_font_settings(
        text_frame.paragraphs[0], font_name, font_size_main, yellow_text, is_bold=True
    )
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
    
    # Check if there's a second line of text for the same slide
    if i + 1 < len(text_chunks):
        # Add a text box for the second line of text
        top = Inches(0.7)  # Position the second line below the first one
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text_chunks[i + 1]  # Second line of text
        apply_font_settings(
            text_frame.paragraphs[0], font_name, font_size_main, yellow_text, is_bold=True
        )
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text

    # Add the last encountered P text to the bottom of the slide if it exists
    if p_text:
        bottom_textbox = slide.shapes.add_textbox(left, Inches(5.2), width, height)
        bottom_text_frame = bottom_textbox.text_frame
        bottom_text_frame.text = p_text[p_text_Num]  # P line text
        apply_font_settings(
            bottom_text_frame.paragraphs[0], font_name, font_size_bottom, yellow_text
        )
        bottom_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text

# Save the PowerPoint presentation
presentation.save(output_path)
print(f"Presentation saved to {output_path}")
