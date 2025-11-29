import os
import re
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# --- 核心功能函式 (Core Functions) ---

def find_ppt_path(root_path, song_title):
    """
    在指定的根目錄下遞迴搜尋符合歌名的PPT檔案路徑。
    """
    for dirpath, _, filenames in os.walk(root_path):
        for filename in filenames:
            if not (filename.lower().endswith(".pptx") or filename.lower().endswith(".ppt")):
                continue

            name_stem = os.path.splitext(filename)[0]
            
            # 從檔名中分離出真正的歌曲標題 (支援 '-' 或 ' ' 分隔)
            title_from_file = ""
            if '-' in name_stem:
                title_from_file = name_stem.split('-', 1)[-1].strip()
            elif ' ' in name_stem:
                title_from_file = name_stem.split(' ', 1)[-1].strip()
            else:
                title_from_file = name_stem.strip()
            
            if title_from_file == song_title:
                return os.path.join(dirpath, filename)
    return None

def clean_text(text):
    """
    清除文字中不相容XML的控制字元，但保留換行、歸位和定位符。
    """
    return re.sub(r'[\x00-\x08\x0e-\x1f]', '', text)

def extract_song_data(song_order, ppt_library_path):
    """
    從原始PPT檔案中提取歌名和歌詞數據。
    返回一個包含歌曲資訊的列表。
    """
    print("--- 階段一：提取歌詞資料 ---")
    songs_data = []
    for song_title in song_order:
        print(f"正在處理歌曲: {song_title}...")
        ppt_path = find_ppt_path(ppt_library_path, song_title)
        
        song_info = {'title': song_title, 'lyrics': [], 'is_ppt_old_format': False, 'found': False}

        if not ppt_path:
            print(f"  > 警告: 在資料庫中找不到 '{song_title}' 的檔案。")
            songs_data.append(song_info)
            continue
        
        song_info['found'] = True

        if ppt_path.lower().endswith(".ppt"):
            print(f"  > 警告: '{song_title}' 是舊版 .ppt 格式，無法自動提取歌詞。")
            song_info['is_ppt_old_format'] = True
            songs_data.append(song_info)
            continue

        print(f"  > 找到了: {os.path.basename(ppt_path)}")
        prs = Presentation(ppt_path)
        
        lyrics = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame or not shape.text.strip():
                    continue

                # 假設歌詞主要位於投影片的上半部 (可根據實際情況調整)
                if shape.top < prs.slide_height / 2:
                    raw_text = shape.text_frame.text
                    text_with_standard_breaks = raw_text.replace('\v', '\n')
                    lines = text_with_standard_breaks.splitlines()
                    
                    for line in lines:
                        cleaned_line = clean_text(line.strip())
                        if cleaned_line:
                            lyrics.append(cleaned_line)
        
        song_info['lyrics'] = lyrics
        songs_data.append(song_info)
        
    print("--- 歌詞資料提取完成 ---")
    return songs_data

def generate_word_document(songs_data, template_path, output_path):
    """
    根據提取的歌曲資料和範本，生成Word大字報文件。
    """
    print("--- 階段二：生成 Word 大字報 ---")
    try:
        document = Document(template_path)
    except Exception as e:
        print(f"錯誤：無法載入Word範本 '{template_path}'。請確認檔案是否存在。")
        print(f"詳細錯誤: {e}")
        return

    # 刪除範本頂部的空段落
    if len(document.paragraphs) > 0 and not document.paragraphs[0].text.strip():
        p = document.paragraphs[0]
        p._p.getparent().remove(p._p)

    for song in songs_data:
        document.add_paragraph(f"【{song['title']}】", style='SongTitle')

        if not song['lyrics']:
            if not song['found']:
                p = document.add_paragraph()
                p.add_run("【警告：在詩歌庫中找不到這首歌的檔案，請檢查歌名是否完全匹配】").bold = True
            elif song['is_ppt_old_format']:
                p = document.add_paragraph()
                p.add_run("【注意：此歌曲為舊版.ppt格式，無法自動匯入，請手動處理】").bold = True
            else:
                 p = document.add_paragraph()
                 p.add_run("【注意：找到了檔案，但未能成功提取任何歌詞】").bold = True
            document.add_paragraph("")
            continue

        for line in song['lyrics']:
            p = document.add_paragraph(line, style='Lyrics')
        
        document.add_paragraph("") # 在每首歌后添加一個空段落

    document.save(output_path)
    print(f"--- Word 大字報已成功儲存至 '{output_path}' ---")

def generate_projection_ppt(songs_data, output_path, lines_per_slide=4):
    """
    根據提取的歌曲資料，生成風格統一的投影用PPT。
    """
    print("--- 階段三：生成投影用 PPT ---")
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)
    blank_slide_layout = presentation.slide_layouts[6] # 使用空白版面配置

    # 樣式設定
    black_fill = PptxRGBColor(0, 0, 0)
    yellow_text = PptxRGBColor(255, 255, 0)
    font_name = "微軟正黑體"
    font_size = Pt(44)
    font_size_title = Pt(24)

    for song in songs_data:
        # 處理歌詞頁
        if not song['lyrics']:
            continue # 如果沒有歌詞，就沒有歌詞頁

        # 將歌詞按指定行數分組
        for i in range(0, len(song['lyrics']), lines_per_slide):
            slide = presentation.slides.add_slide(blank_slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = black_fill

            # 歌詞內容
            lyrics_chunk = song['lyrics'][i:i + lines_per_slide]
            lyrics_text = "\n".join(lyrics_chunk)
            
            # 將歌詞文字框置頂
            left, top, width, height = Inches(1), Inches(0.1), Inches(14), Inches(7.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = lyrics_text
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP # 垂直置頂

            for para in text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                para.font.name = font_name
                para.font.size = font_size
                para.font.bold = True
                para.font.color.rgb = yellow_text
            
            # 頁尾歌名置底
            footer_box = slide.shapes.add_textbox(Inches(1), Inches(8.25), Inches(14), Inches(0.75))
            footer_frame = footer_box.text_frame
            footer_frame.text = song['title']
            p_footer = footer_frame.paragraphs[0]
            p_footer.alignment = PP_ALIGN.CENTER
            p_footer.font.name = font_name
            p_footer.font.size = font_size_title
            p_footer.font.color.rgb = yellow_text

    presentation.save(output_path)
    print(f"--- 投影用 PPT 已成功儲存至 '{output_path}' ---")


def main():
    """
    主執行函式
    """
    # 階段一: 從原始PPTs提取所有歌曲資料
    songs_data = extract_song_data(SONG_ORDER, PPT_LIBRARY_PATH)

    # 階段二: 使用提取的資料生成Word文件
    generate_word_document(songs_data, TEMPLATE_DOCX_PATH, OUTPUT_DOCX_PATH)

    # 階段三: 使用提取的資料生成投影PPT
    # 您可以在這裡調整每張投影片的歌詞行數，例如 lines_per_slide=2 或 4
    generate_projection_ppt(songs_data, OUTPUT_PPTX_PATH, lines_per_slide=2)
    
    print("\n所有任務已完成！")

# --- 程式入口 ---
if __name__ == "__main__":
    main()